#!/usr/bin/env python3
from __future__ import annotations

# Standard Library
from pathlib import Path
import shutil
import subprocess
import tempfile

# PIP3 modules
try:
	from pptx import Presentation
except Exception:
	Presentation = None

try:
	from odf import text as odf_text
	from odf.opendocument import load as odf_load
	from odf import teletype as odf_teletype
except Exception:
	odf_text = None
	odf_load = None
	odf_teletype = None

# local repo modules
from .base import FileMetadata, FileMetadataPlugin
from .mdls_utils import mdls_field

#============================================


class PresentationPlugin(FileMetadataPlugin):
	"""
	Plugin for presentation files.
	"""

	name = "presentation"
	supported_suffixes: set[str] = {"ppt", "pptx", "odp"}

	#============================================
	def supports(self, path: Path) -> bool:
		return path.suffix.lower().lstrip(".") in self.supported_suffixes

	#============================================
	def extract_metadata(self, path: Path) -> FileMetadata:
		meta = FileMetadata(path=path, plugin_name=self.name)
		meta.extra["size_bytes"] = path.stat().st_size
		meta.extra["extension"] = path.suffix.lstrip(".")
		title = mdls_field(path, "kMDItemTitle")
		if title:
			meta.title = title
		summary = self._read_preview(path)
		if summary:
			meta.summary = summary
		else:
			meta.summary = f"Presentation file {path.name}"
		return meta

	#============================================
	def _print_why(self, path: Path, message: str) -> None:
		print(f"[WHY] {path.name}: {message}")

	#============================================
	def _read_preview(self, path: Path) -> str | None:
		ext = path.suffix.lower().lstrip(".")
		if ext == "pptx":
			return self._read_pptx_preview(path)
		if ext == "ppt":
			return self._read_ppt_via_soffice(path)
		if ext == "odp" and odf_load and odf_text and odf_teletype:
			try:
				document = odf_load(str(path))
			except Exception:
				return None
			paras = document.getElementsByType(odf_text.P)
			snippets: list[str] = []
			for para in paras:
				text = odf_teletype.extractText(para).strip()
				if text:
					snippets.append(text)
				if len(snippets) >= 40:
					break
			if not snippets:
				return None
			full = " ".join(snippets)
			flat = " ".join(full.split())
			return flat[:1200] if flat else None
		return None

	def _read_pptx_preview(self, path: Path) -> str | None:
		if not Presentation:
			return None
		try:
			prs = Presentation(str(path))
		except Exception:
			return None
		text_runs: list[str] = []
		for slide in prs.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					text = str(shape.text).strip()
					if text:
						text_runs.append(text)
			if len(text_runs) >= 20:
				break
		if not text_runs:
			return None
		full = " ".join(text_runs)
		flat = " ".join(full.split())
		return flat[:1200] if flat else None

	def _read_ppt_via_soffice(self, path: Path) -> str | None:
		if not Presentation:
			self._print_why(path, "python-pptx not installed; skipping PPT conversion")
			return None
		soffice = shutil.which("soffice")
		if not soffice:
			self._print_why(path, "LibreOffice not found; skipping PPT conversion")
			return None
		with tempfile.TemporaryDirectory() as tmp_dir:
			output_path = Path(tmp_dir) / f"{path.stem}.pptx"
			try:
				subprocess.run(
					[
						soffice,
						"--headless",
						"--convert-to",
						"pptx",
						"--outdir",
						tmp_dir,
						str(path),
					],
					check=True,
					stdout=subprocess.DEVNULL,
					stderr=subprocess.DEVNULL,
					timeout=30,
				)
			except Exception as exc:
				self._print_why(path, f"LibreOffice conversion failed ({exc.__class__.__name__})")
				return None
			if not output_path.exists():
				self._print_why(path, "LibreOffice conversion produced no output")
				return None
			return self._read_pptx_preview(output_path)
